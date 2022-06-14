import os

from pptx import Presentation
import pytest

from utils.navigation import paths_keeper
from utils.comparators import compare_texts


def extract_text(prs):
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    yield run.text


@pytest.mark.parametrize(
    "baseline_path,subject_path,result_path",
    paths_keeper.get_paths("pptx"),
    ids=paths_keeper.get_ids("pptx"),
)
def test_pptx_text(baseline_path, subject_path, result_path):
    baseline_presentation = Presentation(baseline_path)
    subject_presentation = Presentation(subject_path)
    compare_texts(
        "\n".join(extract_text(baseline_presentation)),
        "\n".join(extract_text(subject_presentation)),
        os.path.join(result_path, "text_result.txt"),
    )
