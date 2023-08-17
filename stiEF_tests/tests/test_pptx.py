import os

from pptx import Presentation

from stiEF_tests.tests.utils.comparators import compare_texts


def extract_text(prs):
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    yield run.text



def test_pptx_text(reference_path, subject_path, result_path, naming_prefix):
    baseline_presentation = Presentation(reference_path)
    subject_presentation = Presentation(subject_path)
    compare_texts(
        "\n".join(extract_text(baseline_presentation)),
        "\n".join(extract_text(subject_presentation)),
        os.path.join(result_path, "text_result.txt"),
    )
